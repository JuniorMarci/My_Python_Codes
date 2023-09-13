import pandas as pd
import os
#import seaborn
import matplotlib.pyplot as plt
from Funcoes import converterPNUM, Agrupar, GerGrafico, Filtrar, Odf, NumPMes

#print(os.listdir('..\..\\Correções (V2)\\Saída_Algoritmo'))

TabMetas = r"..\..\Correções (V2)\Saída_Algoritmo\Metas.csv"
TabMetas = pd.read_csv(TabMetas,sep=";", thousands='.')

converterPNUM(TabMetas,'Previsto Pontual')
converterPNUM(TabMetas,'Realizado Pontual')
converterPNUM(TabMetas,'Nota Pontual')
converterPNUM(TabMetas,'Nota Acumulado')
converterPNUM(TabMetas,'Realizado Acumulado')
converterPNUM(TabMetas,'Previsto Acumulado')

Categorias = r"..\..\Correções (V2)\Saída_Algoritmo\Categorias.csv"
Categorias = pd.read_csv(Categorias,sep=";", names=['Cod.Unico','Shopping'],skiprows=1)

TabMetas['Departamento'] = [str(meta).split('_')[0] for meta in TabMetas['Código da Meta']]
TabMetas = TabMetas.merge(Categorias, on='Cod.Unico', how='right')#.to_csv("Right.csv",sep=';')


#colors = ['r' if value < 90 else 'y' if value < 99 else 'g' for value in GRUPO['Nota Acumulado']]

#TabMetas['Nota Acumulado'] = TabMetas['Nota Acumulado'].str.replace(',', '.').astype(float)


def GerarApresentacao(TabFiltrada,Dep, mes):
    TabMetas = TabFiltrada

    Od = Odf(-1)
    Tdf = []
    GRUPODep = Filtrar(TabMetas,'',True,'','') # Filtrar(TabMetas,mes,ona,ColFil,ColIg)

    GRUPODep['IdTemp'] = GRUPODep['Departamento']+'--'+GRUPODep['mês'].astype('int').astype('str')
    GRUPODepA = Agrupar(GRUPODep,'IdTemp',"Nota Acumulado")
    GRUPODepA['Departamento'] = GRUPODepA.reset_index()['IdTemp'].str.split('--',expand=True)[0].values
    GRUPODepA['mês'] = GRUPODepA.reset_index()['IdTemp'].str.split('--',expand=True)[1].values.astype('int')
    GRUPODepA['mesT'] = NumPMes(GRUPODepA['mês'])
    GRUPODepA = GRUPODepA.sort_values(by='mês')
    #GRUPODepA.to_csv('T.csv',sep=';')

    #GRUPODep = Agrupar(GRUPODep,'Departamento',"Nota Acumulado") # Agrupar(Tab,Coluna,Metrica)
    GerGrafico(GRUPODepA,'mesT','Nota Acumulado',"",'Departamento '+Dep,Od) # GerGrafico(df,ex,ey,ez,Cat,Od)
    Od = Odf(Od)
    Tdf.append(GRUPODepA)

    for c in GRUPODep['Shopping'].unique():
        GRUPOSh = Filtrar(TabMetas,'',True,'Shopping',c) # Filtrar(TabMetas,mes,ona,ColFil,ColIg)
        GRUPOSh['IdTemp'] = GRUPOSh['Shopping']+'--'+GRUPOSh['mês'].astype('int').astype('str')
        GRUPOShA = Agrupar(GRUPOSh,'IdTemp',"Nota Acumulado")
        GRUPOShA['Shopping'] = GRUPOShA.reset_index()['IdTemp'].str.split('--',expand=True)[0].values
        GRUPOShA['mês'] = GRUPOShA.reset_index()['IdTemp'].str.split('--',expand=True)[1].values.astype('int')
        GRUPOShA['mesT'] = NumPMes(GRUPOShA['mês'])
        GRUPOShA = GRUPOShA.sort_values(by='mês')

        GerGrafico(GRUPOShA,'mesT','Nota Acumulado',"",'Departamento '+Dep+', Unidade '+c,Od) # GerGrafico(df,ex,ey,ez,Cat,Od)
        Od = Odf(Od)
        Tdf.append(GRUPOShA)

        #GRUPOShA.to_csv('c.csv',sep=';')

        for d in GRUPOSh['Login'].unique():

            # Desempenho do Colaborador por mês

            GRUPOCol = Filtrar(TabMetas,'',True,'Login',d) # Filtrar(TabMetas,mes,ona,ColFil,ColIg)
            GRUPOCol['IdTemp'] = GRUPOCol['Login']+'--'+GRUPOCol['mês'].astype('int').astype('str')
            GRUPOColA = Agrupar(GRUPOCol,'IdTemp',"Nota Acumulado")
            GRUPOColA['Login'] = GRUPOColA.reset_index()['IdTemp'].str.split('--',expand=True)[0].values
            GRUPOColA['mês'] = GRUPOColA.reset_index()['IdTemp'].str.split('--',expand=True)[1].values.astype('int')
            GRUPOColA['mesT'] = NumPMes(GRUPOColA['mês'])
            GRUPOColA = GRUPOColA.sort_values(by='mês')

            GerGrafico(GRUPOColA,'mesT','Nota Acumulado',"",'Departamento '+Dep+', Unidade '+c+', Colaborador '+d,Od) # GerGrafico(df,ex,ey,ez,Cat,Od)
            Od = Odf(Od)
            Tdf.append(GRUPOColA)

            #GRUPOColA.to_csv('d.csv',sep=';')

            # Desempenho do Colaborador por metas

            GRUPOCol = Filtrar(TabMetas,mes,True,'Login',d) # Filtrar(TabMetas,mes,ona,ColFil,ColIg)
            GRUPOColA = Agrupar(GRUPOCol,'Código da Meta',"Nota Acumulado")

            GerGrafico(GRUPOColA.reset_index(),'Código da Meta','Nota Acumulado',"",'Metas do colaborador '+d+' no mês '+str(mes),Od) # GerGrafico(df,ex,ey,ez,Cat,Od)

            Od = Odf(Od)
            Tdf.append(Tdf[len(Tdf)-1])

            #GRUPOColA.to_csv('dd.csv',sep=';')

            for e in GRUPOCol['Código da Meta'].unique():

                GRUPOMes = Filtrar(TabMetas,"",False,'Código da Meta',e) # Filtrar(TabMetas,mes,ona,ColFil,ColIg)

                GRUPOMes['mêsT'] = NumPMes(GRUPOMes['mês'])
                GerGrafico(GRUPOMes,'mêsT','Realizado Pontual','Previsto Pontual','Colaborador '+d+', Meta '+e,Od) # GerGrafico(df,ex,ey,ez,Cat,Od)
                Od = Odf(Od)
                Tdf.append(GRUPOMes)


    '''
    for c in GRUPODep.index: # Iterando sobre cada departamento, nota dos shoppings
        print(c)
        GRUPOSh = Filtrar(TabMetas,'',True,'Departamento',c) # Filtrar(TabMetas,mes,ona,ColFil,ColIg)
        GRUPOShA = Agrupar(GRUPOSh,'Shopping',"Nota Acumulado") # Agrupar(Tab,Coluna,Metrica)
        if (len(GRUPOShA.index) > 1): GerGrafico(GRUPOShA,'Shopping','Nota Acumulado',"",'Departamento '+c,Od) # GerGrafico(df,ex,ey,ez,Cat,Od)
        Od = Odf(Od)
        Tdf.append(GRUPOSh)


        for d in GRUPOShA.index: # Iterando sobre cada shopping, nota dos colaboradores
            print(d)
            #GRUPOCol = Filtrar(GRUPOSh,mes,True,'Shopping',d) # Filtrar(TabMetas,mes,ona,ColFil,ColIg)

            GRUPOCol = Filtrar(GRUPOSh,'',True,'Shopping',d)
            GRUPOCol['IdTemp'] = GRUPOCol['Login']+'--'+GRUPOCol['mês'].astype('int').astype('str')
            GRUPOColA = Agrupar(GRUPOCol,'IdTemp',"Nota Acumulado")
            GRUPOColA['Login'] = GRUPOColA.reset_index()['IdTemp'].str.split('--',expand=True)[0].values
            GRUPOColA['mês'] = GRUPOColA.reset_index()['IdTemp'].str.split('--',expand=True)[1].values.astype('int')
            GRUPOColA['mesT'] = NumPMes(GRUPOColA['mês'])
            GRUPOColA = GRUPOColA.sort_values(by='mês')
            GRUPOColA.to_csv(d+'.csv',sep=';')

            #GRUPOColA = Agrupar(GRUPOCol,'Login',"Nota Acumulado") # Agrupar(Tab,Coluna,Metrica)
            #GerGrafico(GRUPOColA,'Login','Nota Acumulado',"",'Departamento '+c+", Unidade "+d,Od) # GerGrafico(df,ex,ey,ez,Cat,Od)

            GerGrafico(GRUPOColA,'mesT','Nota Acumulado',"",'Departamento '+c+", Unidade "+d+' -- Testes',Od)
            Od = Odf(Od)
            Tdf.append(GRUPOCol)

            for e in GRUPOColA.index: # Iterando sobre cada colaborador, nota por meta
                print(e)
                GRUPOMet = Filtrar(GRUPOCol,mes,True,'Login',e) # Filtrar(TabMetas,mes,ona,ColFil,ColIg)
                GRUPOMetA = Agrupar(GRUPOMet,'Código da Meta',"Nota Acumulado") # Agrupar(Tab,Coluna,Metrica)
                GerGrafico(GRUPOMetA,'Código da Meta','Nota Acumulado',"",'Departamento '+c+", Unidade "+d+", Colaborador "+e,Od) # GerGrafico(df,ex,ey,ez,Cat,Od)
                Od = Odf(Od)
                Tdf.append(GRUPOMet)

                for f in GRUPOMetA.index: # Iterando sobre cada meta, nota por mês
                    print(f)
                    GRUPOMes = Filtrar(TabMetas,"",False,'Código da Meta',f) # Filtrar(TabMetas,mes,ona,ColFil,ColIg)

                    GRUPOMes['mêsT'] = NumPMes(GRUPOMes['mês'])
                    GerGrafico(GRUPOMes,'mêsT','Realizado Pontual','Previsto Pontual','Colaborador '+e+', Meta '+f,Od) # GerGrafico(df,ex,ey,ez,Cat,Od)
                    Od = Odf(Od)
                    Tdf.append(GRUPOMes)

    '''



    # ---------------------------------------


    import pptx

    # Criar um objeto Presentation vazio
    prs = pptx.Presentation()

    # Definir o tamanho da apresentação para A4 paisagem
    e = 2
    w = 29.7
    h = 21
    prs.slide_width = pptx.util.Cm(w*e)
    prs.slide_height = pptx.util.Cm(h*e)

    def NovoSlide(ima,txt):

        # Adicionar um slide em branco
        blank_slide_layout = prs.slide_master.slide_layouts[6]

        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        # Configurações do slide e do título
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(ima, pptx.util.Cm(4.3), 0)
        Titulo = slide.shapes.add_textbox(50, 50, pptx.util.Cm(w*e), pptx.util.Cm(20)).text_frame
        Titulo.text = txt
        Titulo.fit_text(font_family='Calibri', max_size=30, bold=False, italic=False, font_file=None)
        Titulo.paragraphs[0].alignment = PP_ALIGN.CENTER

        # add_table(rows, cols, left, top, width, height)
        Linhas = 9
        Colunas = 13
        Tabela = slide.shapes.add_table(Linhas, Colunas, 0, pptx.util.Cm(23), pptx.util.Cm(w*e), pptx.util.Cm(15))

        #Tabela.table.cell(0,0).text = "Oi"

        #Colocar toda a tabela com fundo branco
        for row in Tabela.table.rows:
            for cell in row.cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

        def InsTxt(Tab,x,y,txt):
            if (txt=="-1.0"): txt = "N/A"
            Tab.table.cell(x, y).text_frame.text = txt
            if ('Valor' in txt): Tab.table.cell(x, y).text_frame.paragraphs[0].runs[0].font.size = pptx.util.Pt(30)  # Change the font size to 30 points
            elif (len(txt)<6): Tab.table.cell(x, y).text_frame.paragraphs[0].runs[0].font.size = pptx.util.Pt(30)  # Change the font size to 30 points
            else: Tab.table.cell(x, y).text_frame.paragraphs[0].runs[0].font.size = pptx.util.Pt(20)  # Change the font size to 30 points
            Tab.table.cell(x, y).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            Tab.table.cell(x, y).text_frame.paragraphs[0].font.bold = True
            if ('Valor' in txt):
                Tab.table.cell(x, y).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
                Tab.table.cell(x, y).fill.solid()
                Tab.table.cell(x, y).fill.fore_color.rgb = RGBColor(210,4,4)

            else:
                Tab.table.cell(x, y).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                Tab.table.cell(x, y).fill.solid()
                #FORECO = Tab.table.cell(x, y).fill.fore_color
                #FORECO.rgb = RGBColor(255,255,255)
                Tab.table.cell(x, y).fill.fore_color.rgb = RGBColor(255,255,255)

        Meses = ['Jan','Fev','Mar','Abr','Mai','Jun','Jul','Ago','Set','Out','Nov','Dez']
        for c in Meses:
            InsTxt(Tabela,1,Meses.index(c)+1,c)

        # Mesclagem das células de título
        InsTxt(Tabela,0,0,'Valor Pontual')
        Tabela.table.rows[0].cells[0].merge(Tabela.table.rows[0].cells[12]) # Merge the cells in the first row
        InsTxt(Tabela,5,0,'Valor Acumulado')
        Tabela.table.rows[5].cells[0].merge(Tabela.table.rows[5].cells[12])

        InsTxt(Tabela,2,0,'P')
        InsTxt(Tabela,3,0,'R')
        InsTxt(Tabela,4,0,'D')
        InsTxt(Tabela,6,0,'P')
        InsTxt(Tabela,7,0,'R')
        InsTxt(Tabela,8,0,'D')

        # Identificação do banco/filtro a ser usado com base na ordem (var Od) que fica carimbada nas iniciais do nome da imagem
        PosAt = int(txt.split('-')[0])
        PosAt = Tdf[PosAt]
        TipoAg = txt.split('-')[1].replace(' ','')

        def DEB():
            try:
                print('Começo')
                print(PosAt['Nota Acumulado'])
                print(PosAt['mês'])
            except Exception as EXX: print(EXX)

        if (TipoAg == 'RealizadoPontual'):
            try:
                for ms in range(len(Meses)):
                    MA = (PosAt.reset_index()['mês'][ms]) # Valor do mês, linkado com a posição reiniciada
                    #print(PosAt.reset_index()['Previsto Pontual'][ms])
                    InsTxt(Tabela, 2, int(MA), str(PosAt.reset_index()['Previsto Pontual'][ms])) # Valor Pontual, linkado com a posição reiniciada
                    InsTxt(Tabela, 3, int(MA), str(PosAt.reset_index()['Realizado Pontual'][ms])) # Valor Pontual, linkado com a posição reiniciada
                    InsTxt(Tabela, 6, int(MA), str(PosAt.reset_index()['Previsto Acumulado'][ms])) # Valor Pontual, linkado com a posição reiniciada
                    InsTxt(Tabela, 7, int(MA), str(PosAt.reset_index()['Realizado Acumulado'][ms])) # Valor Pontual, linkado com a posição reiniciada
                    DP = float(PosAt.reset_index()['Realizado Pontual'][ms])-float(PosAt.reset_index()['Previsto Pontual'][ms])
                    DA = float(PosAt.reset_index()['Realizado Acumulado'][ms])-float(PosAt.reset_index()['Previsto Acumulado'][ms])
                    InsTxt(Tabela,4, int(MA), str(int(DP)))
                    InsTxt(Tabela,8, int(MA), str(int(DA)))

            except: 1==1 #print('Sem mês')

        if (TipoAg == 'NotaAcumulado'):
            #DEB()
            try:
                for ms in range(len(Meses)):
                    MA = (PosAt.reset_index()['mês'][ms]) # Valor do mês, linkado com a posição reiniciada
                    #print(PosAt.reset_index()['Previsto Pontual'][ms])
                    #InsTxt(Tabela, 2, int(MA), str(PosAt.reset_index()['Previsto Pontual'][ms])) # Valor Pontual, linkado com a posição reiniciada
                    InsTxt(Tabela, 3, int(MA), str(PosAt.reset_index()['Nota Pontual'][ms])) # Valor Pontual, linkado com a posição reiniciada
                    #InsTxt(Tabela, 6, int(MA), str(PosAt.reset_index()['Previsto Acumulado'][ms])) # Valor Pontual, linkado com a posição reiniciada
                    InsTxt(Tabela, 7, int(MA), str(PosAt.reset_index()['Nota Acumulado'][ms])) # Valor Pontual, linkado com a posição reiniciada
                    #DP = float(PosAt.reset_index()['Realizado Pontual'][ms])-float(PosAt.reset_index()['Previsto Pontual'][ms])
                    #DA = float(PosAt.reset_index()['Realizado Acumulado'][ms])-float(PosAt.reset_index()['Previsto Acumulado'][ms])
                    #InsTxt(Tabela,4, int(MA), str(int(DP)))
                    #InsTxt(Tabela,8, int(MA), str(int(DA)))

            except Exception as err: print('Sem mês',err)


    for c in os.listdir('../Graficos'):
        NovoSlide('../Graficos/'+c,c)
        print('Slide ',c)
        os.remove('../Graficos/'+c)

    # Salvar a apresentação
    prs.save('../Apresentacao__'+Dep+'.pptx')



#'''
for Dep in TabMetas[~ pd.isna(TabMetas['Departamento'])]['Departamento'].unique():
    print(Dep)
    TabM = Filtrar(TabMetas,"",False,'Departamento',Dep)
    GerarApresentacao(TabM,Dep,7)

'''
TabM = Filtrar(TabMetas,"",False,'Departamento','DFIN')
GerarApresentacao(TabM,'DFIN',7)
'''